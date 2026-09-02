import os
import streamlit as st
from dotenv import load_dotenv
from google import genai
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity


# ============================================================
# CONFIGURATION
# ============================================================

load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

GEMINI_MODEL = "gemini-3.6-flash"


# ============================================================
# PAGE CONFIGURATION
# ============================================================

st.set_page_config(
    page_title="TeachMate AI",
    page_icon="👩‍🏫",
    layout="wide"
)


# ============================================================
# GEMINI CONNECTION
# ============================================================

if not GEMINI_API_KEY:
    st.error(
        "GEMINI_API_KEY is missing. "
        "Please check your .env file."
    )
    st.stop()

client = genai.Client(
    api_key=GEMINI_API_KEY
)


# ============================================================
# SESSION STATE
# ============================================================

if "lesson" not in st.session_state:
    st.session_state.lesson = ""

if "quiz" not in st.session_state:
    st.session_state.quiz = ""

if "evaluation" not in st.session_state:
    st.session_state.evaluation = ""

if "document_text" not in st.session_state:
    st.session_state.document_text = ""

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []


# ============================================================
# GEMINI FUNCTION
# ============================================================

def ask_gemini(prompt):

    try:

        interaction = client.interactions.create(
            model=GEMINI_MODEL,
            input=prompt,
            generation_config={
                "thinking_level": "low"
            }
        )

        return interaction.output_text

    except Exception as e:

        return f"ERROR: {str(e)}"


# ============================================================
# PDF EXTRACTION
# ============================================================

def extract_pdf(file):

    text = ""

    reader = PdfReader(file)

    for page in reader.pages:

        page_text = page.extract_text()

        if page_text:
            text += page_text + "\n"

    return text


# ============================================================
# DOCX EXTRACTION
# ============================================================

def extract_docx(file):

    document = Document(file)

    text = []

    for paragraph in document.paragraphs:

        if paragraph.text.strip():

            text.append(paragraph.text)

    return "\n".join(text)


# ============================================================
# PPTX EXTRACTION
# ============================================================

def extract_pptx(file):

    presentation = Presentation(file)

    text = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                if shape.text.strip():

                    slide_text.append(shape.text)

        if slide_text:

            text.append(
                f"Slide {slide_number}\n"
                + "\n".join(slide_text)
            )

    return "\n\n".join(text)


# ============================================================
# TXT EXTRACTION
# ============================================================

def extract_txt(file):

    return file.read().decode(
        "utf-8",
        errors="ignore"
    )


# ============================================================
# DOCUMENT PROCESSING
# ============================================================

def extract_document(file):

    filename = file.name.lower()

    if filename.endswith(".pdf"):

        return extract_pdf(file)

    elif filename.endswith(".docx"):

        return extract_docx(file)

    elif filename.endswith(".pptx"):

        return extract_pptx(file)

    elif filename.endswith(".txt"):

        return extract_txt(file)

    return ""


# ============================================================
# CREATE TEXT CHUNKS
# ============================================================

def create_chunks(
    text,
    chunk_size=800,
    overlap=100
):

    words = text.split()

    chunks = []

    start = 0

    while start < len(words):

        end = min(
            start + chunk_size,
            len(words)
        )

        chunk = " ".join(
            words[start:end]
        )

        if chunk.strip():

            chunks.append(chunk)

        if end >= len(words):

            break

        start = end - overlap

    return chunks


# ============================================================
# RAG RETRIEVAL
# ============================================================

def retrieve_context(
    document_text,
    query,
    top_k=5
):

    if not document_text.strip():

        return ""

    chunks = create_chunks(
        document_text
    )

    if not chunks:

        return ""

    if len(chunks) == 1:

        return chunks[0]

    try:

        vectorizer = TfidfVectorizer(
            stop_words="english"
        )

        matrix = vectorizer.fit_transform(
            chunks + [query]
        )

        similarities = cosine_similarity(
            matrix[-1],
            matrix[:-1]
        ).flatten()

        top_indices = (
            similarities
            .argsort()[-top_k:][::-1]
        )

        selected = [
            chunks[i]
            for i in top_indices
        ]

        return "\n\n".join(selected)

    except Exception:

        return "\n\n".join(
            chunks[:top_k]
        )


# ============================================================
# GENERATE PERSONALIZED LESSON
# ============================================================

def generate_lesson(
    topic,
    level,
    language,
    learning_style,
    context
):

    prompt = f"""
You are TeachMate AI, a personalized AI teacher.

Student level:
{level}

Preferred language:
{language}

Learning style:
{learning_style}

Topic:
{topic}

Relevant uploaded study material:
{context if context else "No study material uploaded."}

Teach this topic clearly and accurately.

Follow these requirements:

1. Start with a simple introduction.
2. Explain the concept step-by-step.
3. Adapt the explanation to the student's level.
4. Adapt the explanation to the student's learning style.
5. Give simple real-world examples.
6. Use an analogy where useful.
7. Highlight important points.
8. Include a short "Key Takeaways" section.
9. Include three "Check Your Understanding" questions.
10. If uploaded material is provided, use it as the primary reference.
11. Do not invent information from the uploaded material.

Write the lesson in the student's preferred language.
"""

    return ask_gemini(prompt)


# ============================================================
# GENERATE VIDEO SCRIPT
# ============================================================

def generate_video_script(
    topic,
    lesson
):

    prompt = f"""
Create a natural teaching script for an AI teacher.

Topic:
{topic}

Lesson:
{lesson}

The script should:

- Sound like a real teacher speaking.
- Be friendly and engaging.
- Explain the main concept simply.
- Give one real-world example.
- Avoid markdown.
- Avoid mentioning AI, Gemini, APIs or software.
- Be suitable for an avatar video.
- Be approximately 1 minute long.

Return only the spoken script.
"""

    return ask_gemini(prompt)


# ============================================================
# GENERATE QUIZ
# ============================================================

def generate_quiz(
    topic,
    lesson
):

    prompt = f"""
Create a 5-question assessment.

Topic:
{topic}

Lesson:
{lesson}

For each question provide:

Question:
A.
B.
C.
D.
Correct Answer:
Explanation:

Questions should test actual understanding,
not just memorization.
"""

    return ask_gemini(prompt)


# ============================================================
# EVALUATE STUDENT UNDERSTANDING
# ============================================================

def evaluate_student(
    topic,
    lesson,
    student_answer
):

    prompt = f"""
You are an adaptive AI teacher.

Topic:
{topic}

Lesson:
{lesson}

Student answer:
{student_answer}

Analyze the student's understanding.

Provide:

1. Overall correctness
2. What the student understood
3. What the student misunderstood
4. Possible misconception
5. A simpler explanation
6. A new example
7. One follow-up question

Be encouraging and supportive.
"""

    return ask_gemini(prompt)


# ============================================================
# HEADER
# ============================================================

st.title("👩‍🏫 TeachMate AI")

st.markdown(
    """
### Your Personalized AI Teaching Assistant

Learn • Ask • Practice • Get Feedback • Improve
"""
)


# ============================================================
# SIDEBAR
# ============================================================

with st.sidebar:

    st.header("🎓 Student Profile")

    student_name = st.text_input(
        "Student Name",
        "Student"
    )

    level = st.selectbox(
        "Learning Level",
        [
            "Beginner",
            "Intermediate",
            "Advanced"
        ]
    )

    language = st.selectbox(
        "Preferred Language",
        [
            "English",
            "Hindi",
            "Telugu",
            "Tamil",
            "Kannada"
        ]
    )

    learning_style = st.selectbox(
        "Learning Style",
        [
            "Simple explanations",
            "Examples first",
            "Step-by-step",
            "Visual learning"
        ]
    )

    st.divider()

    st.info(
        "TeachMate AI adapts explanations "
        "and assessments to the learner."
    )


# ============================================================
# TOPIC AND MATERIAL
# ============================================================

st.header("📚 1. Choose Your Topic")

topic = st.text_input(
    "What do you want to learn?",
    placeholder="Example: Newton's Laws of Motion"
)

uploaded_file = st.file_uploader(
    "Upload study material (optional)",
    type=[
        "pdf",
        "docx",
        "pptx",
        "txt"
    ]
)


# ============================================================
# PROCESS UPLOADED FILE
# ============================================================

if uploaded_file:

    try:

        document_text = extract_document(
            uploaded_file
        )

        st.session_state.document_text = (
            document_text
        )

        if document_text.strip():

            st.success(
                f"✅ {uploaded_file.name} "
                "processed successfully."
            )

            with st.expander(
                "📄 Preview uploaded material"
            ):

                st.write(
                    document_text[:5000]
                )

        else:

            st.warning(
                "The uploaded file contains "
                "no readable text."
            )

    except Exception as e:

        st.error(
            f"File processing error: {e}"
        )


# ============================================================
# LESSON GENERATION
# ============================================================

if st.button(
    "🧠 Generate Personalized Lesson",
    type="primary",
    use_container_width=True
):

    if not topic.strip():

        st.warning(
            "Please enter a topic first."
        )

    else:

        with st.spinner(
            "TeachMate AI is preparing your lesson..."
        ):

            context = retrieve_context(
                st.session_state.document_text,
                topic
            )

            lesson = generate_lesson(
                topic,
                level,
                language,
                learning_style,
                context
            )

            if lesson.startswith("ERROR:"):

                st.error(lesson)

            else:

                st.session_state.lesson = lesson

                st.session_state.quiz = ""

                st.session_state.evaluation = ""

                st.success(
                    "🎉 Your personalized lesson is ready!"
                )


# ============================================================
# DISPLAY LESSON
# ============================================================

if st.session_state.lesson:

    st.divider()

    st.header("👩‍🏫 2. Personalized Lesson")

    st.markdown(
        st.session_state.lesson
    )


    # ========================================================
    # AI TEACHER VIDEO
    # ========================================================

    st.divider()

    st.header("🎥 3. AI Teacher Video")

    st.info(
        "Your human-like AI teacher video "
        "was created using HeyGen."
    )

    st.write(
        "For the hackathon demonstration, "
        "the HeyGen avatar video can be shown "
        "as a separate demonstration because "
        "your current HeyGen API trial credits "
        "have been exhausted."
    )

    st.caption(
        "AI Teacher: HeyGen Avatar Demonstration"
    )


    # ========================================================
    # ASK AI TEACHER
    # ========================================================

    st.divider()

    st.header("💬 4. Ask Your AI Teacher")

    question = st.text_input(
        "Ask anything about this topic",
        placeholder="Example: Can you explain this with a simple example?"
    )

    if st.button(
        "💡 Ask Teacher",
        use_container_width=True
    ):

        if question.strip():

            context = retrieve_context(
                st.session_state.document_text,
                question
            )

            prompt = f"""
You are TeachMate AI.

Student level:
{level}

Learning style:
{learning_style}

Preferred language:
{language}

Topic:
{topic}

Lesson:
{st.session_state.lesson}

Relevant uploaded material:
{context}

Student question:
{question}

Answer clearly and patiently.

Adapt the answer to the student's level.

Use an example or analogy when useful.

If uploaded material is available,
prioritize information from it.
"""

            with st.spinner(
                "Your AI teacher is answering..."
            ):

                answer = ask_gemini(
                    prompt
                )

            if answer.startswith("ERROR:"):

                st.error(answer)

            else:

                st.session_state.chat_history.append(
                    {
                        "question": question,
                        "answer": answer
                    }
                )

    for chat in st.session_state.chat_history:

        st.markdown(
            f"**🧑‍🎓 You:** {chat['question']}"
        )

        st.markdown(
            f"**👩‍🏫 TeachMate AI:** {chat['answer']}"
        )


    # ========================================================
    # ADAPTIVE LEARNING
    # ========================================================

    st.divider()

    st.header("🧠 5. Check My Understanding")

    student_answer = st.text_area(
        "Explain what you learned in your own words.",
        placeholder="Write your answer here..."
    )

    if st.button(
        "🔍 Evaluate My Understanding",
        use_container_width=True
    ):

        if not student_answer.strip():

            st.warning(
                "Please write an answer first."
            )

        else:

            with st.spinner(
                "Analyzing your understanding..."
            ):

                evaluation = evaluate_student(
                    topic,
                    st.session_state.lesson,
                    student_answer
                )

                st.session_state.evaluation = (
                    evaluation
                )

    if st.session_state.evaluation:

        st.markdown(
            st.session_state.evaluation
        )


    # ========================================================
    # QUIZ
    # ========================================================

    st.divider()

    st.header("📝 6. Final Assessment")

    if st.button(
        "📝 Generate 5-Question Quiz",
        use_container_width=True
    ):

        with st.spinner(
            "Creating your assessment..."
        ):

            quiz = generate_quiz(
                topic,
                st.session_state.lesson
            )

            st.session_state.quiz = quiz

    if st.session_state.quiz:

        st.markdown(
            st.session_state.quiz
        )


# ============================================================
# FOOTER
# ============================================================

st.divider()

st.caption(
    "TeachMate AI | Personalized AI Teaching Assistant"
)